import fitz  # PyMuPDF
import requests
import io
import os
import tempfile
import logging
import pandas as pd
import pytesseract
from docx import Document
from PIL import Image
from pptx import Presentation
from typing import List, Optional

logger = logging.getLogger(__name__)

class DocumentParserService:
    @staticmethod
    def parse_pdf(path_or_url: str) -> List[str]:
        """Parses a PDF and returns a list of per-page strings."""
        try:
            if path_or_url.startswith("http"):
                response = requests.get(path_or_url)
                response.raise_for_status()
                file_bytes = io.BytesIO(response.content)
                doc = fitz.open(stream=file_bytes, filetype="pdf")
            else:
                doc = fitz.open(path_or_url)

            pages = [page.get_text().strip() for page in doc]
            doc.close()
            return [p for p in pages if p]
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            return []

    @staticmethod
    def parse_docx(path: str) -> List[str]:
        """Parses a DOCX and returns a list of chunks."""
        try:
            doc = Document(path)
            text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
            
            # Use ~2000 chars as a "page" equivalent
            chunk_size = 2000
            return [text[i:i + chunk_size].strip() for i in range(0, len(text), chunk_size) if text[i:i + chunk_size].strip()]
        except Exception as e:
            logger.error(f"Error parsing DOCX: {e}")
            return []

    @staticmethod
    def parse_image(path: str) -> List[str]:
        """Extracts text from image using OCR and return as list."""
        try:
            image = Image.open(path)
            raw_text = pytesseract.image_to_string(image).strip()
            return [raw_text] if raw_text else []
        except Exception as e:
            logger.error(f"Error parsing image: {e}")
            return []

    @staticmethod
    def parse_pptx(path: str) -> List[str]:
        """Parses a PPTX and returns a list of per-slide strings."""
        try:
            prs = Presentation(path)
            slides_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    elif shape.shape_type == 19:  # Table
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text: slide_text.append(" | ".join(row_text))
                    elif shape.shape_type == 13:  # Picture
                         try:
                            # Attempt OCR on slide images if possible
                            image_part = shape.image.blob
                            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                                tmp.write(image_part)
                                tmp_path = tmp.name
                            ocr_text = pytesseract.image_to_string(Image.open(tmp_path)).strip()
                            if ocr_text: slide_text.append(f"[OCR]: {ocr_text}")
                            os.unlink(tmp_path)
                         except:
                             pass
                if slide_text:
                    slides_content.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
            return slides_content
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            return []

    @staticmethod
    def parse_excel(path: str) -> List[str]:
        """Parses Excel and returns string representation per sheet."""
        try:
            xls = pd.ExcelFile(path)
            content = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                if not df.empty:
                    content.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}")
            return content
        except Exception as e:
            logger.error(f"Error parsing Excel: {e}")
            return []

    @classmethod
    def parse_any(cls, path: str) -> List[str]:
        """Auto-detects format and parses."""
        ext = path.lower().split('.')[-1]
        if ext == 'pdf': return cls.parse_pdf(path)
        if ext in ['docx', 'doc']: return cls.parse_docx(path)
        if ext in ['pptx', 'ppt']: return cls.parse_pptx(path)
        if ext in ['xlsx', 'xls', 'csv']: return cls.parse_excel(path)
        if ext in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']: return cls.parse_image(path)
        return []
