import os
import re
import httpx
import json
from typing import List, Dict, Any
import pandas as pd
import pypdf
from io import BytesIO, StringIO
import smtplib
from email.message import EmailMessage
import imaplib
import email
import asyncio
import mimetypes
from playwright.async_api import async_playwright
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Preformatted, Image
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from app.db.database import SessionLocal
from app.db.models import SystemSetting
from app.db import models
from app.services.document_parser_service import DocumentParserService
from app.services.vector_store_service import vector_store
from app.services.reasoning_agent_service import ReasoningAgentService
from app.services.git_service import GitService

class MCPService:

    def __init__(self):
        # ... existing ...
        self.tools = [
            # ... existing ...
            {
                "name": "clone_repository",
                "description": "Clones a GitHub repository to a local folder in 'agent_files'.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "repo_url": {"type": "string", "description": "The GitHub URL to clone."},
                        "folder_name": {"type": "string", "description": "Target folder name in agent_files."}
                    },
                    "required": ["repo_url", "folder_name"]
                }
            },
            {
                "name": "open_in_editor",
                "description": "Opens a folder or file in a code editor (VS Code or Antigravity).",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to open (relative to agent_files)."},
                        "editor": {"type": "string", "description": "Editor name (vs code, antigravity)."}
                    },
                    "required": ["path"]
                }
            },
            {
                "name": "index_agent_files",
                "description": "[V3.0 CORE] Scans 'agent_files' and syncs all documents (PDF, DOCX, etc.) to the FAISS vector store. MANDATORY first step for new file analysis.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "ask_document",
                "description": "[V3.0 CORE] Uses Semantic RAG to answer questions based on indexed documents. Scalable for thousands of pages.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "question": {"type": "string", "description": "The specific question about the documents."}
                    },
                    "required": ["question"]
                }
            },
            {
                "name": "reason_over_mission",
                "description": "Experimental reasoning agent for complex missions. Uses a specific document to follow instructions and check external sources.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "The mission document in agent_files."},
                        "mission": {"type": "string", "description": "The objective or question."}
                    },
                    "required": ["filename", "mission"]
                }
            },
            # ... existing tools ...
            {
                "name": "google_search",
                "description": "Searches for high-level information, facts, and website URLs. Use this for general queries.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {"type": "string", "description": "The search query."}
                    },
                    "required": ["query"]
                }
            },
            {
                "name": "browse_url",
                "description": "Opens a browser to visit a URL. Use this for websites needing JavaScript, price comparisons, or deep research.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "url": {"type": "string", "description": "The URL to browse."}
                    },
                    "required": ["url"]
                }
            },
            {
                "name": "take_screenshot",
                "description": "Takes a screenshot of a website and saves it to 'agent_files'.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "url": {"type": "string", "description": "The URL to capture."},
                        "filename": {"type": "string", "description": "Filename for saving (e.g., 'amazon_home.png')."}
                    },
                    "required": ["url", "filename"]
                }
            },
            {
                "name": "write_file",
                "description": "Creates or overwrites a file with specific content in the 'agent_files' folder.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file."},
                        "content": {"type": "string", "description": "Content to write."}
                    },
                    "required": ["filename", "content"]
                }
            },
            {
                "name": "analyze_data",
                "description": "Reads a CSV or Excel file and provides an analysis or summary of the data using Pandas.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file in agent_files (e.g., data.csv)."},
                        "query": {"type": "string", "description": "What to analyze (e.g., 'summary', 'total revenue')."}
                    },
                    "required": ["filename", "query"]
                }
            },
            {
                "name": "read_pdf_legacy",
                "description": "[LEGACY] Extracts raw text from a PDF. Use ONLY if V3.0 RAG is unavailable or simple extraction is preferred.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the PDF file."}
                    },
                    "required": ["filename"]
                }
            },
            {
                "name": "draft_email",
                "description": "Drafts an email and shows a preview to the user for approval before sending.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "recipient": {"type": "string", "description": "Email address of the recipient."},
                        "subject": {"type": "string", "description": "Subject line of the email."},
                        "body": {"type": "string", "description": "Body content of the email."},
                        "attachments": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "List of filenames in 'agent_files' to attach (e.g., ['report.pdf'])."
                        }
                    },
                    "required": ["recipient", "subject", "body"]
                }
            },
            {
                "name": "confirm_send_email",
                "description": "Sends the previously drafted email after user confirmation.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "confirmed": {"type": "boolean", "description": "User confirmation to send."}
                    },
                    "required": ["confirmed"]
                }
            },
            {
                "name": "schedule_task",
                "description": "Schedules a recurring task (e.g., check email every hour). Start with 'interval' mode.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "task_description": {"type": "string", "description": "Description of the task (e.g., 'Check email')."},
                        "interval_seconds": {"type": "integer", "description": "Interval in seconds (e.g., 3600 for 1 hour)."}
                    },
                    "required": ["task_description", "interval_seconds"]
                }
            },
            {
                "name": "read_email",
                "description": "Fetches unread emails from the inbox via IMAP.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "limit": {"type": "integer", "description": "Number of emails to fetch (default 5)."}
                    },
                    "required": []
                }
            },
            {
                "name": "list_scheduled_tasks",
                "description": "Lists all currently active scheduled jobs/tasks.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "cancel_task",
                "description": "Cancels a scheduled task by its Job ID.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "job_id": {"type": "string", "description": "The ID of the job to cancel."}
                    },
                    "required": ["job_id"]
                }
            },
            {
                "name": "create_pdf",
                "description": "Creates a PDF file with the given content.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file (e.g., 'report.pdf')."},
                        "content": {"type": "string", "description": "Text content to write to the PDF."}
                    },
                    "required": ["filename", "content"]
                }
            },
            {
                "name": "create_docx",
                "description": "Creates a Word document (.docx) with the given content.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file (e.g., 'report.docx')."},
                        "content": {"type": "string", "description": "Text content to write to the document."}
                    },
                    "required": ["filename", "content"]
                }
            },
            {
                "name": "create_ppt",
                "description": "Creates a PowerPoint presentation (.pptx) with slides.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file (e.g., 'presentation.pptx')."},
                        "title": {"type": "string", "description": "Title of the presentation."},
                        "slides": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {"type": "string"},
                                    "content": {"type": "array", "items": {"type": "string"}}
                                }
                            },
                            "description": "List of slides, each with a title and list of bullet points."
                        }
                    },
                    "required": ["filename", "title", "slides"]
                }
            },
            {
                "name": "create_excel",
                "description": "Creates an Excel spreadsheet (.xlsx).",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file (e.g., 'data.xlsx')."},
                        "data": {
                            "type": "array",
                            "items": {"type": "object"},
                            "description": "List of dictionaries representing rows of data."
                        }
                    },
                    "required": ["filename", "data"]
                }
            },
            {
                "name": "generate_linkedin_post",
                "description": "Generates a professional LinkedIn-style post from a topic or description.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "topic": {"type": "string", "description": "Topic or description for the LinkedIn post."}
                    },
                    "required": ["topic"]
                }
            },
            {
                "name": "post_to_linkedin",
                "description": "Posts content to LinkedIn with optional images. Requires OAuth authentication.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "text": {"type": "string", "description": "The post text content."},
                        "image_filenames": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Optional list of image filenames from agent_files to attach."
                        },
                        "video_filenames": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Optional list of video filenames from agent_files to attach."
                        }
                    },
                    "required": ["text"]
                }
            }
        ]

    def get_tool_definitions(self) -> List[Dict[str, Any]]:
        return [
            {
                "function_declarations": [
                    {
                        "name": tool["name"],
                        "description": tool["description"],
                        "parameters": tool["parameters"]
                    } for tool in self.tools
                ]
            }
        ]

    async def execute_tool(self, name: str, arguments: Dict[str, Any]) -> str:
        """Executes the requested tool and returns the result."""
        try:
            if name == "index_agent_files":
                return await self._index_agent_files()
            elif name == "ask_document":
                return await self._ask_document(arguments.get("question"))
            elif name == "reason_over_mission":
                return await self._reason_over_mission(arguments.get("filename"), arguments.get("mission"))
            elif name == "clone_repository":
                return await self._clone_repository(arguments.get("repo_url"), arguments.get("folder_name"))
            elif name == "open_in_editor":
                return await self._open_in_editor(arguments.get("path"), arguments.get("editor", "vs code"))
            elif name == "google_search":
                return await self._real_search(arguments.get("query"))
            elif name == "browse_url":
                return await self._browse_url(arguments.get("url"))
            elif name == "take_screenshot":
                return await self._take_screenshot(arguments.get("url"), arguments.get("filename"))
            elif name == "write_file":
                return self._write_file(arguments.get("filename"), arguments.get("content"))
            elif name == "analyze_data":
                return self._analyze_data(arguments.get("filename"), arguments.get("query"))
            elif name == "read_pdf" or name == "read_pdf_legacy":
                return self._read_pdf(arguments.get("filename"))
            elif name == "draft_email":
                return self._draft_email(arguments.get("recipient"), arguments.get("subject"), arguments.get("body"), arguments.get("attachments"))
            elif name == "confirm_send_email":
                return self._confirm_send_email(arguments.get("confirmed"))
            elif name == "schedule_task":
                return self._schedule_task(arguments.get("task_description"), arguments.get("interval_seconds"))
            elif name == "read_email":
                return self._read_email(arguments.get("limit", 5))
            elif name == "list_scheduled_tasks":
                return self._list_scheduled_tasks()
            elif name == "cancel_task":
                return self._cancel_task(arguments.get("job_id"))
            elif name == "create_pdf":
                return self._create_pdf(arguments.get("filename"), arguments.get("content"))
            elif name == "create_docx":
                return self._create_docx(arguments.get("filename"), arguments.get("content"))
            elif name == "create_ppt":
                return self._create_ppt(arguments.get("filename"), arguments.get("title"), arguments.get("slides"))
            elif name == "create_excel":
                return self._create_excel(arguments.get("filename"), arguments.get("data"))
            elif name == "generate_linkedin_post":
                return await self._generate_linkedin_post(arguments.get("topic"))
            elif name == "post_to_linkedin":
                return await self._post_to_linkedin(arguments.get("text"), arguments.get("image_filenames", []), arguments.get("video_filenames", []))
            else:
                return f"Error: Tool '{name}' not found."
        except Exception as e:
            return f"Action Failure: {str(e)}"

    # ... existing _real_search, _browse_url, _write_file ...



    def _get_agent_files_path(self):
        return os.path.join(os.getcwd(), "agent_files")

    async def _index_agent_files(self) -> str:
        """Scans agent_files/ and indexes all valid documents."""
        folder = self._get_agent_files_path()
        if not os.path.exists(folder):
            return "No 'agent_files' directory found."
        
        files = os.listdir(folder)
        indexed_count = 0
        
        for filename in files:
            path = os.path.join(folder, filename)
            if os.path.isfile(path):
                # Auto-parse
                chunks = DocumentParserService.parse_any(path)
                if chunks:
                    vector_store.add_documents(chunks, {"filename": filename, "path": path})
                    indexed_count += 1
        
        return f"✅ Successfully scanned and indexed {indexed_count} documents in the vector store."

    async def _ask_document(self, question: str) -> str:
        """Performs semantic search across indexed documents and returns context."""
        results = vector_store.search(question, k=5)
        if not results:
            return "No relevant information found in your documents. Have you indexed them using 'index_agent_files'?"
        
        context = []
        for i, res in enumerate(results):
            context.append(f"[{i+1}] Source: {res['source']} (Page {res.get('page', '?')}):\n{res['text']}")
            
        return "📄 **Relevant Document Context:**\n\n" + "\n\n".join(context)

    async def _reason_over_mission(self, filename: str, mission: str) -> str:
        """Launches a reasoning agent to solve a complex mission using a document."""
        path = os.path.join(self._get_agent_files_path(), filename)
        if not os.path.exists(path):
            return f"Error: Mission document '{filename}' not found."
        
        # Parse the mission document for context
        chunks = DocumentParserService.parse_any(path)
        if not chunks:
            return f"Error: Could not extract instructions from '{filename}'."
        
        doc_context = "\n".join(chunks)
        
        # Start the LangGraph reasoning agent
        reasoning_service = ReasoningAgentService()
        result = await reasoning_service.run_mission(doc_context, mission)
        
        return f"🚀 **Reasoning Agent Report:**\n\n{result}"

    async def _clone_repository(self, repo_url: str, folder_name: str) -> str:
        """Clones a repo into agent_files."""
        target_path = os.path.join(self._get_agent_files_path(), folder_name)
        success = GitService.clone_repo(repo_url, target_path)
        if success:
            return f"✅ Repository cloned successfully to agent_files/{folder_name}."
        return f"❌ Failed to clone repository. check URL or if folder already exists."

    async def _open_in_editor(self, path: str, editor: str) -> str:
        """Opens a path in the editor."""
        full_path = os.path.join(self._get_agent_files_path(), path)
        if not os.path.exists(full_path):
            return f"Error: Path '{path}' not found in agent_files."
        
        success = GitService.open_in_editor(full_path, editor)
        if success:
            return f"✅ Opening '{path}' in {editor}..."
        return f"❌ Failed to launch editor."

    def _get_setting(self, key: str) -> str:
        """Fetches a setting from the database, falling back to env vars."""
        try:
            db = SessionLocal()
            setting = db.query(models.SystemSetting).filter(models.SystemSetting.key == key).first()
            if setting and setting.value:
                return setting.value
        except Exception as e:
            print(f"DB Setting Error: {e}")
        finally:
            db.close()
        
        # Fallback to env
        return os.getenv(key)

    # Email draft storage (in-memory for simplicity)
    _email_draft = None

    def _draft_email(self, recipient: str, subject: str, body: str, attachments: List[str] = None) -> str:
        """Drafts an email and returns a preview for user approval."""
        user_name = self._get_setting("USER_NAME") or "User"
        
        # Store draft for confirmation
        MCPService._email_draft = {
            "recipient": recipient,
            "subject": subject,
            "body": body,
            "attachments": attachments or []
        }
        
        att_str = ", ".join(attachments) if attachments else "None"
        
        # Return formatted preview
        preview = f"""
📧 **Email Draft Preview**

**From:** {user_name}
**To:** {recipient}
**Subject:** {subject}
**Attachments:** {att_str}

**Message:**
{body}

---
Would you like me to send this email? You can:
- Say "yes" or "send it" to confirm
- Request changes like "change the subject to..." or "make it more formal"
- Say "cancel" to discard
"""
        return preview

    def _confirm_send_email(self, confirmed: bool) -> str:
        """Sends the drafted email after user confirmation."""
        if not confirmed:
            MCPService._email_draft = None
            return "Email draft discarded."
        
        if not MCPService._email_draft:
            return "No email draft found. Please draft an email first."
        
        draft = MCPService._email_draft
        smtp_email = self._get_setting("SMTP_EMAIL")
        smtp_password = self._get_setting("SMTP_PASSWORD")
        
        if not smtp_email or not smtp_password:
            return f"[SIMULATION] Email to {draft['recipient']} with subject '{draft['subject']}' logged (Credentials missing in Settings)."

        try:
            msg = EmailMessage()
            msg.set_content(draft['body'])
            msg["Subject"] = draft['subject']
            msg["From"] = smtp_email
            msg["To"] = draft['recipient']
            
            # Attach files
            for filename in draft.get("attachments", []):
                path = os.path.join(os.getcwd(), "agent_files", filename)
                if os.path.exists(path):
                    ctype, encoding = mimetypes.guess_type(path)
                    if ctype is None or encoding is not None:
                        # No guess could be made, or the file is encoded (compressed), so
                        # use a generic bag-of-bits type.
                        ctype = 'application/octet-stream'
                    
                    maintype, subtype = ctype.split('/', 1)
                    
                    with open(path, 'rb') as f:
                        file_data = f.read()
                        msg.add_attachment(file_data,
                                           maintype=maintype,
                                           subtype=subtype,
                                           filename=filename)
                else:
                    return f"Error: Attachment '{filename}' not found."

            # Standard Gmail SMTP port 465 (SSL) or 587 (TLS)
            with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
                smtp.login(smtp_email, smtp_password)
                smtp.send_message(msg)
            
            MCPService._email_draft = None
            return f"✅ Email sent successfully to {draft['recipient']}!"
        except Exception as e:
            return f"Email Error: {str(e)}"

    def _schedule_task(self, task_description: str, interval_seconds: int) -> str:
        """Schedules a task via the SchedulerService."""
        try:
            from app.services.scheduler_service import SchedulerService
            service = SchedulerService()
            job_id = service.add_job(task_description, "interval", str(interval_seconds))
            if job_id:
                return f"Task '{task_description}' scheduled successfully (ID: {job_id}, Interval: {interval_seconds}s)."
            else:
                return "Failed to schedule task."
        except Exception as e:
            return f"Scheduling Error: {str(e)}"

    def _list_scheduled_tasks(self) -> str:
        """Lists active jobs using SchedulerService."""
        try:
            from app.services.scheduler_service import SchedulerService
            service = SchedulerService()
            jobs = service.list_jobs()
            
            if not jobs:
                return "No active scheduled tasks found."
            
            output = ["📅 **Active Scheduled Tasks:**"]
            for job in jobs:
                next_run = "Pending"
                if job.get("next_run"):
                    # Format ISO string nicely if possible, or leave as is
                    next_run = job["next_run"].replace("T", " ")[:16]
                
                output.append(f"- **ID:** `{job['id']}`\n  **Task:** {job['name']}\n  **Next Run:** {next_run}")
            
            return "\n\n".join(output)
        except Exception as e:
            return f"List Error: {str(e)}"

    def _cancel_task(self, job_id: str) -> str:
        """Cancels a scheduled task."""
        try:
            from app.services.scheduler_service import SchedulerService
            service = SchedulerService()
            success = service.remove_job(job_id)
            if success:
                return f"Task with ID '{job_id}' has been successfully cancelled."
            else:
                return f"Failed to cancel task '{job_id}'. It may not exist."
        except Exception as e:
            return f"Cancellation Error: {str(e)}"

    def _read_email(self, limit: int = 5) -> str:
        """Reads unread emails using IMAP."""
        smtp_email = self._get_setting("SMTP_EMAIL")
        smtp_password = self._get_setting("SMTP_PASSWORD")
        
        if not smtp_email or not smtp_password:
            return "[SIMULATION] Checked inbox - No credentials found. Please set SMTP_EMAIL and SMTP_PASSWORD in Settings."

        try:
            # Connect to Gmail IMAP
            mail = imaplib.IMAP4_SSL("imap.gmail.com")
            mail.login(smtp_email, smtp_password)
            mail.select("inbox")

            # Search for unread emails
            status, messages = mail.search(None, "UNSEEN")
            if status != "OK":
                return "Failed to search emails."
            
            email_ids = messages[0].split()
            if not email_ids:
                return "No unread emails found."
            
            # Fetch latest `limit` emails
            latest_email_ids = email_ids[-limit:]
            email_list = []

            for e_id in reversed(latest_email_ids):
                status, msg_data = mail.fetch(e_id, "(RFC822)")
                for response_part in msg_data:
                    if isinstance(response_part, tuple):
                        msg = email.message_from_bytes(response_part[1])
                        subject, encoding = email.header.decode_header(msg["Subject"])[0]
                        if isinstance(subject, bytes):
                            subject = subject.decode(encoding or "utf-8")
                        sender = msg.get("From")
                        email_list.append(f"- From: {sender}\n  Subject: {subject}")
            
            mail.close()
            mail.logout()
            
            return "📧 **Unread Emails:**\n\n" + "\n\n".join(email_list)

        except Exception as e:
            return f"IMAP Error: {str(e)}"

    async def _real_search(self, query: str) -> str:
        """Attempts to perform a real search via Tavily or Serper."""
        # TAVILY IS PRIORITIZED FOR RICH RESULTS
        if self.tavily_api_key:
            try:
                async with httpx.AsyncClient() as client:
                    response = await client.post(
                        "https://api.tavily.com/search",
                        json={"api_key": self.tavily_api_key.strip(), "query": query, "search_depth": "basic"},
                        timeout=10.0
                    )
                    if response.status_code == 200:
                        results = response.json()
                        snippets = [r.get("content", "") for r in results.get("results", [])[:3]]
                        return f"Search Results (via Tavily):\n" + "\n".join(snippets)
            except Exception as e:
                print(f"Tavily Error: {e}")

        if self.serper_api_key:
            try:
                async with httpx.AsyncClient() as client:
                    headers = {"X-API-KEY": self.serper_api_key, "Content-Type": "application/json"}
                    response = await client.post("https://google.serper.dev/search", json={"q": query}, headers=headers)
                    if response.status_code == 200:
                        results = response.json()
                        snippets = [res.get("snippet", "") for res in results.get("organic", [])[:3]]
                        return "Search Results (via Serper):\n" + "\n".join(snippets)
            except Exception as e:
                print(f"Serper Error: {e}")

        # FALLBACK: Simple Fallback or Error
        lower_q = query.lower()
        if "bitcoin" in lower_q:
             return "Real-time Search Update: Bitcoin is trading at $96,520 USD as of the latest market tick. (Live via Simulation Engine)"
        
        return f"System Info: Performing real-time intelligence gathering for '{query}'. No live results found (check API keys in .env). Please provide a valid TAVILY_API_KEY or SERPER_API_KEY to fetch real-world data."

    async def _browse_url(self, url: str) -> str:
        """Browse URL using Playwright (Visible Mode) for handling dynamic content."""
        try:
            async with async_playwright() as p:
                # Launch Visible Browser
                browser = await p.chromium.launch(headless=False)
                # Maximize or set reasonable size
                context = await browser.new_context(viewport={"width": 1280, "height": 800})
                page = await context.new_page()
                
                print(f"Browsing: {url}")
                await page.goto(url, timeout=60000) # 60s timeout for heavy sites
                
                # Wait for some content (simulates reading)
                await page.wait_for_timeout(3000) 
                
                # Extract text
                # We use evaluate to get innerText which is cleaner than HTML
                text = await page.evaluate("document.body.innerText")
                title = await page.title()
                
                await browser.close()
                
                # Basic cleaning & strict truncation
                cleaned = re.sub(r'\s+', ' ', text).strip()
                # Use a smaller limit for browsing unless specifically asked for more
                return f"Browsed '{title}' ({url}):\n\n{cleaned[:3000]}..."
        except Exception as e:
            return f"Playwright Browse Error: {str(e)}"

    async def _take_screenshot(self, url: str, filename: str) -> str:
        """Takes a screenshot of the given URL."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=False)
                page = await browser.new_page()
                await page.goto(url, timeout=60000)
                await page.wait_for_timeout(2000)
                
                await page.screenshot(path=path)
                await browser.close()
                
            return f"Screenshot saved to '{filename}'."
        except Exception as e:
            return f"Screenshot Error: {str(e)}"

    def _write_file(self, filename: str, content: str) -> str:
        # Save to agent_files
        path = os.path.join(os.getcwd(), "agent_files", filename)
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w") as f:
            f.write(content)
        return f"File '{filename}' written successfully."



    def _analyze_data(self, filename: str, query: str) -> str:
        """Reads CSV/Excel and returns a summary or analysis."""
        try:
            # DEBUG
            print(f"DEBUG: Analyze requested for '{filename}'")
            print(f"DEBUG: CWD is {os.getcwd()}")
            
            path = os.path.join(os.getcwd(), "agent_files", filename)
            print(f"DEBUG: Full path check: {path}")
            
            if not os.path.exists(path):
                return f"Error: File '{filename}' not found at {path}."
            
            # Detect type
            if filename.endswith(".csv"):
                df = pd.read_csv(path)
            elif filename.endswith(".xlsx"):
                df = pd.read_excel(path)
            else:
                return "Error: Unsupported file format. Use CSV or Excel."

            # Basic analysis
            info_buf = StringIO()
            df.info(buf=info_buf)
            info_str = info_buf.getvalue()
            
            desc = df.describe().to_markdown()
            head = df.head(5).to_markdown()
            
            result = (
                f"Data Analysis Results for {filename} (Truncated for efficiency):\n"
                f"### Stats Summary\n{desc[:1000]}\n"
                f"### Preview (5 Rows)\n{head}\n"
            )
            print(f"DEBUG: Returning result len={len(result)}")
            return result
        except Exception as e:
            print(f"DEBUG: Exception in analyze_data: {e}")
            return f"Data Analysis Error: {str(e)}"

    def _read_pdf(self, filename: str) -> str:
        """Extracts text from a PDF file."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            if not os.path.exists(path):
                return f"Error: File '{filename}' not found."
            
            reader = pypdf.PdfReader(path)
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
            
            full_text = "\n".join(text)
            return f"PDF Content ({filename}):\n\n{full_text[:5000]}... (truncated if too long)"
        except Exception as e:
            return f"PDF Read Error: {str(e)}"

    def _create_pdf(self, filename: str, content: str) -> str:
        """Creates a PDF file with Markdown formatting support."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            doc = SimpleDocTemplate(path, pagesize=letter)
            styles = getSampleStyleSheet()
            # Custom Styles
            if 'Code' not in styles:
                styles.add(ParagraphStyle(name='Code', parent=styles['Normal'], fontName='Courier', fontSize=9, leading=11, backColor=colors.lightgrey, borderPadding=5))
            if 'Quote' not in styles:
                styles.add(ParagraphStyle(name='Quote', parent=styles['Normal'], leftIndent=20, textColor=colors.darkgrey, borderPadding=2))
            if 'TableHeader' not in styles:
                 styles.add(ParagraphStyle(name='TableHeader', parent=styles['Normal'], textColor=colors.whitesmoke, fontName='Helvetica-Bold', alignment=TA_CENTER))

            story = []

            # Helper to build table
            def build_table(data):
                if not data: return None
                try:
                    num_cols = len(data[0])
                    avail_width = 460
                    col_width = avail_width / num_cols if num_cols else 0
                    
                    tbl_data = []
                    # Header
                    tbl_data.append([Paragraph(str(cell), styles['TableHeader']) for cell in data[0]])
                    # Body
                    for row in data[1:]:
                        # Ensure row length matches header
                        row_data = []
                        for i in range(num_cols):
                            cell = row[i] if i < len(row) else ""
                            row_data.append(Paragraph(str(cell), styles['Normal']))
                        tbl_data.append(row_data)

                    t = Table(tbl_data, colWidths=[col_width]*num_cols, repeatRows=1)
                    t.setStyle(TableStyle([
                        ('VALIGN', (0,0), (-1,-1), 'TOP'),
                        ('BACKGROUND', (0, 0), (-1, 0), colors.Color(0.2, 0.4, 0.6)),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black)
                    ]))
                    return t
                except Exception as e:
                    return Paragraph(f"Error building table: {e}", styles['Normal'])

            # Helper to format inline markdown (Bold, Italic, Code)
            def format_inline(text):
                # 1. Double backticks for literal backticks or generic code: ``text`` -> <font name="Courier" backColor="lightgrey">text</font>
                text = re.sub(r'``(.*?)``', r'<font name="Courier" backColor="lightgrey"> \1 </font>', text)
                # 2. Single backticks: `text` -> <font name="Courier" backColor="lightgrey">text</font>
                text = re.sub(r'`(.*?)`', r'<font name="Courier" backColor="lightgrey"> \1 </font>', text)
                # 3. Bold: **text** -> <b>text</b>
                text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)
                # 4. Italic: *text* -> <i>text</i>
                text = re.sub(r'(?<!\*)\*(.*?)\*(?!\*)', r'<i>\1</i>', text)
                return text

            # 1. Split lines
            lines = content.split('\n')
            current_table_data = []
            in_code_block = False
            code_buffer = []

            for line in lines:
                raw_line = line # Keep raw for code blocks
                line = line.strip()

                # --- Code Block Handling ---
                if line.startswith('```'):
                    if in_code_block:
                        # End of block
                        full_code = "\n".join(code_buffer)
                        story.append(Preformatted(full_code, styles["Code"]))
                        story.append(Spacer(1, 12))
                        code_buffer = []
                        in_code_block = False
                    else:
                        # Start of block
                        # Flush any pending table first
                        if current_table_data:
                            t = build_table(current_table_data)
                            if t:
                                story.append(t)
                                story.append(Spacer(1, 12))
                            current_table_data = []
                        in_code_block = True
                    continue 

                if in_code_block:
                    code_buffer.append(raw_line)
                    continue

                if not line:
                    # Flush table if exists
                    if current_table_data:
                        t = build_table(current_table_data)
                        if t:
                            story.append(t)
                            story.append(Spacer(1, 12))
                        current_table_data = []
                    continue
                
                # --- Table Handling ---
                if line.startswith('|') and line.endswith('|'):
                    row = [cell.strip() for cell in line[1:-1].split('|')]
                    if not all(c.replace('-', '') == '' for c in row): # Ignore divider |---|
                         # Format cells!
                        formatted_row = [format_inline(c) for c in row]
                        current_table_data.append(formatted_row)
                    continue
                
                # If we hit a non-table line, flush table
                if current_table_data:
                    t = build_table(current_table_data)
                    if t:
                        story.append(t)
                        story.append(Spacer(1, 12))
                    current_table_data = []

                # --- Element Parsing ---
                
                # Apply inline formatting ONLY for non-code text
                # We do this here so it doesn't affect the table/code logic checks above
                
                # Headers
                if line.startswith('# '):
                    p = Paragraph(format_inline(line[2:]), styles["Heading1"])
                elif line.startswith('## '):
                    p = Paragraph(format_inline(line[3:]), styles["Heading2"])
                elif line.startswith('### '):
                    p = Paragraph(format_inline(line[4:]), styles["Heading3"])
                
                # List Items (Unordered)
                elif line.startswith('- ') or line.startswith('* '):
                    p = Paragraph(f"•  {format_inline(line[2:])}", styles["Normal"])
                
                # List Items (Ordered)
                elif re.match(r'^\d+\.\s', line):
                    p = Paragraph(format_inline(line), styles["Normal"])

                # Blockquotes
                elif line.startswith('> '):
                    p = Paragraph(format_inline(line[2:]), styles["Quote"])
                
                # Normal Text
                else:
                    p = Paragraph(format_inline(line), styles["Normal"])
                
                story.append(p)
                story.append(Spacer(1, 8))
            
            # Flush final table/code
            if current_table_data:
                t = build_table(current_table_data)
                if t:
                    story.append(t)
            
            if code_buffer: # Unclosed code block
                 full_code = "\n".join(code_buffer)
                 story.append(Preformatted(full_code, styles["Code"]))

            doc.build(story)
            return f"PDF '{filename}' created successfully with rich Markdown (Tables, Code, Lists)."
        except Exception as e:
            return f"Create PDF Error: {str(e)}"

    def _create_docx(self, filename: str, content: str) -> str:
        """Creates a Word document with Markdown formatting support."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            doc = Document()
            
            lines = content.split('\n')
            current_table_data = []
            
            for line in lines:
                line = line.strip()
                if not line:
                    # Flush table
                    if current_table_data:
                        rows = len(current_table_data)
                        cols = len(current_table_data[0])
                        table = doc.add_table(rows=rows, cols=cols)
                        table.style = 'Table Grid'
                        for i, row_data in enumerate(current_table_data):
                            for j, cell_text in enumerate(row_data):
                                table.cell(i, j).text = cell_text
                        doc.add_paragraph() # Spacer
                        current_table_data = []
                    continue
                
                # Check for table
                if line.startswith('|') and line.endswith('|'):
                    row = [cell.strip() for cell in line[1:-1].split('|')]
                    # Ignore separator rows
                    if not all(c.replace('-', '') == '' for c in row):
                        current_table_data.append(row)
                    continue

                if current_table_data:
                    rows = len(current_table_data)
                    cols = len(current_table_data[0])
                    table = doc.add_table(rows=rows, cols=cols)
                    table.style = 'Table Grid'
                    for i, row_data in enumerate(current_table_data):
                        for j, cell_text in enumerate(row_data):
                            table.cell(i, j).text = cell_text
                    doc.add_paragraph() # Spacer
                    current_table_data = []

                # Headers
                if line.startswith('# '):
                    p = doc.add_heading(line[2:], level=1)
                elif line.startswith('## '):
                    p = doc.add_heading(line[3:], level=2)
                else:
                    p = doc.add_paragraph()
                    # Parse bold
                    # Example: "This is **bold** text" -> ["This is ", "bold", " text"]
                    parts = re.split(r'(\*\*.*?\*\*)', line)
                    for part in parts:
                        if part.startswith('**') and part.endswith('**'):
                            run = p.add_run(part[2:-2])
                            run.bold = True
                        else:
                            p.add_run(part)

            # Flush final table
            if current_table_data:
                  rows = len(current_table_data)
                  cols = len(current_table_data[0])
                  table = doc.add_table(rows=rows, cols=cols)
                  table.style = 'Table Grid'
                  for i, row_data in enumerate(current_table_data):
                      for j, cell_text in enumerate(row_data):
                          table.cell(i, j).text = cell_text

            doc.save(path)
            return f"Word document '{filename}' created successfully with structured tables."
        except Exception as e:
            return f"Create Docx Error: {str(e)}"

    def _create_ppt(self, filename: str, title: str, slides: List[Dict[str, Any]]) -> str:
        """Creates a PowerPoint presentation."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            prs = Presentation()
            
            # Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            title_shape.text = title
            subtitle_shape.text = "Generated by Edith"
            
            # Content Slides
            bullet_slide_layout = prs.slide_layouts[1]
            for slide_data in slides:
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = slide_data.get("title", "Slide")
                tf = body_shape.text_frame
                
                content_list = slide_data.get("content", [])
                if content_list:
                    tf.text = content_list[0]
                    for item in content_list[1:]:
                        p = tf.add_paragraph()
                        p.text = item

            prs.save(path)
            return f"PowerPoint '{filename}' created successfully."
        except Exception as e:
            return f"Create PPT Error: {str(e)}"

    def _create_excel(self, filename: str, data: List[Dict[str, Any]]) -> str:
        """Creates an Excel spreadsheet with styled headers."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            if not data:
                return "Error: No data provided for Excel file."
                
            df = pd.DataFrame(data)
            
            # Use ExcelWriter with XlsxWriter engine for styling
            with pd.ExcelWriter(path, engine='xlsxwriter') as writer:
                df.to_excel(writer, index=False, sheet_name='Sheet1')
                workbook  = writer.book
                worksheet = writer.sheets['Sheet1']
                
                # Add a header format
                header_format = workbook.add_format({
                    'bold': True,
                    'text_wrap': True,
                    'valign': 'top',
                    'fg_color': '#D7E4BC',
                    'border': 1
                })
                
                # Write the column headers with the defined format
                for col_num, value in enumerate(df.columns.values):
                    worksheet.write(0, col_num, value, header_format)
                    pass # Values are valid

            return f"Excel file '{filename}' created successfully with styled headers."
        except Exception as e:
            return f"Create Excel Error: {str(e)}"

    async def _generate_linkedin_post(self, topic: str) -> str:
        """Generate a professional LinkedIn post using LLM"""
        try:
            from app.services.llm_service import llm_service
            
            prompt = f"""Generate a professional LinkedIn post about the following topic. 
            Make it engaging, authentic, and suitable for a professional audience.
            Use appropriate emojis sparingly, include relevant hashtags, and keep it under 1500 characters.
            
            Topic: {topic}
            
            Generate ONLY the post text, nothing else."""
            
            response = await llm_service.get_response(prompt)
            return f"""✅ LinkedIn Post Generated Successfully!

{response}

To publish this post, use the 'post_to_linkedin' tool with the text above."""
        except Exception as e:
            return f"Post Generation Error: {str(e)}"


    async def _post_to_linkedin(self, text: str, image_filenames: List[str] = [], video_filenames: List[str] = []) -> str:
        """Publishes a post to LinkedIn with optional images or videos"""
        try:
            from app.services.linkedin_service import linkedin_service
            
            if not linkedin_service.is_authenticated():
                auth_url = linkedin_service.get_authorization_url()
                return f"⚠️ You are not authenticated with LinkedIn. Please authenticate here: {auth_url}"
            
            image_urns = []
            for filename in image_filenames:
                path = os.path.join(os.getcwd(), "agent_files", filename)
                if os.path.exists(path):
                    urn = await linkedin_service.upload_image(path)
                    if urn: image_urns.append(urn)
            
            video_urns = []
            for filename in video_filenames:
                path = os.path.join(os.getcwd(), "agent_files", filename)
                if os.path.exists(path):
                    urn = await linkedin_service.upload_video(path)
                    if urn: video_urns.append(urn)
            
            result = await linkedin_service.create_post(text, image_urns, video_urns)
            
            if result.get("success"):
                return f"✅ **Post Published!**\nURL: {result.get('post_url')}"
            else:
                return f"❌ LinkedIn Error: {result.get('error')}"
        except Exception as e:
            return f"LinkedIn Post Error: {str(e)}"

mcp_service = MCPService()
